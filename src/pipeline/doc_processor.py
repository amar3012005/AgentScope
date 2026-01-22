# src/pipeline/doc_processor.py

import hashlib
import json
import logging
import os
import re
import statistics
import time
import unicodedata
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import fitz  # PyMuPDF
import yaml
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

# Configure logging for PDF extraction
logger = logging.getLogger(__name__)

# ============================================================================
# OCR SUPPORT NOT YET IMPLEMENTED
# ============================================================================
# The 'enable_ocr' configuration option exists but OCR functionality has not
# been implemented. Scanned PDFs will be extracted as empty or with minimal
# text. To add OCR support, you would need to:
# 1. Install Tesseract OCR and pytesseract
# 2. Implement OCR detection logic in extract_pdf()
# 3. Add image extraction and OCR processing
# ============================================================================

# ============= PDF EXTRACTION CONSTANTS =============
MIN_GAP_WIDTH = 20  # Minimum gap width for column detection
MIDDLE_GAP_WIDTH = 60  # Width for middle gap analysis

# Pre-compiled regex patterns for performance
REGEX_PATTERNS = {
    "page_number": re.compile(r"^[\-\s]*(\d+|[ivxIVX]+)[\-\s/]*\d*[\-\s]*$"),
    "page_text": re.compile(r"^(Page|Seite|Página|页)\s*\d+", re.IGNORECASE),
    "multiple_spaces": re.compile(r" {2,}"),
    "excessive_newlines": re.compile(r"\n{4,}"),
    "triple_newlines": re.compile(r"\n{3}"),
    "space_before_newline": re.compile(r" +\n"),
    "space_after_newline": re.compile(r"\n +"),
}


# ============= PDF EXTRACTION ENUMS AND DATACLASSES =============
class LayoutType(Enum):
    """Page layout types"""

    SINGLE_PORTRAIT = "single_portrait"
    LANDSCAPE_OR_DOUBLE = "landscape_or_double"


class BlockType(Enum):
    """Text block types for better classification"""

    MAIN_TEXT = "main_text"
    TITLE = "title"
    HEADER = "header"
    FOOTER = "footer"
    INFOBOX = "infobox"
    MARGIN_LEFT = "margin_left"
    MARGIN_RIGHT = "margin_right"
    TABLE = "table"
    CAPTION = "caption"
    FOOTNOTE = "footnote"
    PAGE_NUMBER = "page_number"

    def __lt__(self, other):
        """Make BlockType sortable"""
        if not isinstance(other, BlockType):
            return NotImplemented
        return self.value < other.value


@dataclass(slots=True)
class TextBlock:
    """Enhanced text block with position and type information"""

    text: str
    bbox: Tuple[float, float, float, float]  # (x0, y0, x1, y1)
    page_num: int
    block_type: str = "text"
    confidence: float = 1.0
    column_id: Optional[int] = None
    semantic_type: Optional[BlockType] = None
    font_size: Optional[float] = None

    @property
    def x_center(self) -> float:
        return (self.bbox[0] + self.bbox[2]) / 2

    @property
    def y_center(self) -> float:
        return (self.bbox[1] + self.bbox[3]) / 2

    @property
    def width(self) -> float:
        return self.bbox[2] - self.bbox[0]

    @property
    def height(self) -> float:
        return self.bbox[3] - self.bbox[1]


# ============= PDF COMPLEX LAYOUT EXTRACTOR =============
class PDFComplexLayoutExtractor:
    """
    Advanced PDF text extractor with intelligent layout recognition
    and context-preserving reading order.
    """

    def __init__(self, pdf_path: str, debug: bool = False):
        self.pdf_path = pdf_path
        self.debug = debug

        if debug:
            logging.basicConfig(level=logging.DEBUG)
        else:
            logging.basicConfig(level=logging.INFO)

        self.doc = fitz.open(pdf_path)
        self.page_layouts = []
        self.coordinate_system = "unknown"

        self._analyze_page_layouts()
        self._detect_coordinate_system()

    def _analyze_page_layouts(self):
        """Analyze and classify layout type for each page"""
        for page_num, page in enumerate(self.doc):
            rect = page.rect
            width, height = rect.width, rect.height

            blocks = page.get_text("dict").get("blocks", [])
            text_blocks = [b for b in blocks if b.get("type") == 0]

            layout_type = self._detect_page_layout(width, height, text_blocks)
            self.page_layouts.append(
                {
                    "page_num": page_num,
                    "width": width,
                    "height": height,
                    "layout_type": layout_type,
                    "text_blocks": text_blocks,
                }
            )

            logger.debug(f"Page {page_num + 1}: {layout_type.value} ({width:.0f}x{height:.0f})")

    def _detect_page_layout(self, width: float, height: float, blocks: List) -> LayoutType:
        """Detect page layout type"""
        if width <= height * 1.1:
            return LayoutType.SINGLE_PORTRAIT
        return LayoutType.LANDSCAPE_OR_DOUBLE

    def _detect_coordinate_system(self):
        """Detect if Y coordinates are inverted"""
        if not self.doc or len(self.doc) == 0:
            return

        page = self.doc[0]
        blocks = page.get_text("dict").get("blocks", [])
        page_height = page.rect.height

        for block in blocks:
            if block.get("type") != 0:
                continue

            text = self._extract_block_text(block).lower()
            bbox = block.get("bbox", [0, 0, 0, 0])
            y_pos = bbox[1]

            if any(word in text for word in ["page", "seite", "©", "copyright"]):
                if y_pos < page_height * 0.2:
                    self.coordinate_system = "inverted"
                    logger.debug("Detected inverted Y-coordinate system")
                    return
                elif y_pos > page_height * 0.8:
                    self.coordinate_system = "normal"
                    logger.debug("Detected normal Y-coordinate system")
                    return

        self.coordinate_system = "normal"

    def _get_sort_y(self, block: TextBlock) -> float:
        """Get Y coordinate for sorting, accounting for coordinate system"""
        if self.coordinate_system == "inverted":
            page_height = self.page_layouts[block.page_num]["height"]
            return page_height - block.bbox[1]
        else:
            return block.bbox[1]

    def _analyze_middle_gap_coverage(
        self, blocks: List[TextBlock], mid_x: float, gap_width: float = MIDDLE_GAP_WIDTH
    ) -> float:
        """Analyze how much of the middle area is empty (0.0 to 1.0)"""
        gap_left = mid_x - gap_width / 2
        gap_right = mid_x + gap_width / 2

        coverage_points = 0
        total_points = 0

        for y in range(0, int(self.page_layouts[blocks[0].page_num]["height"]), 10):
            total_points += 1
            for block in blocks:
                if (
                    block.bbox[0] <= gap_right
                    and block.bbox[2] >= gap_left
                    and block.bbox[1] <= y <= block.bbox[3]
                ):
                    coverage_points += 1
                    break

        if total_points == 0:
            return 1.0

        empty_ratio = 1.0 - (coverage_points / total_points)
        return empty_ratio

    def _detect_page_numbers(
        self, blocks: List[TextBlock], page_width: float, page_height: float
    ) -> Dict:
        """Detect page numbers and their positions"""
        page_number_info = {"left": False, "right": False, "center": False, "numbers": []}

        for block in blocks:
            if block.width < 100 and block.height < 50:
                if block.y_center < page_height * 0.15 or block.y_center > page_height * 0.85:
                    text = block.text.strip()
                    if REGEX_PATTERNS["page_number"].match(text) or REGEX_PATTERNS[
                        "page_text"
                    ].match(text):
                        page_number_info["numbers"].append(text)

                        if block.x_center < page_width * 0.2:
                            page_number_info["left"] = True
                        elif block.x_center > page_width * 0.8:
                            page_number_info["right"] = True
                        elif 0.4 * page_width < block.x_center < 0.6 * page_width:
                            page_number_info["center"] = True

        return page_number_info

    def _extract_block_text(self, block: Dict) -> str:
        """Extract text from a block dictionary"""
        text_parts = []
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                text = span.get("text", "")
                if not isinstance(text, str):
                    text = str(text)
                text_parts.append(text)
        return " ".join(text_parts)

    def _fix_ligatures(self, text: str) -> str:
        """Replace ligatures with their ASCII equivalents"""
        ligature_map = {
            "\ufb01": "fi",
            "\ufb02": "fl",
            "\ufb00": "ff",
            "\ufb03": "ffi",
            "\ufb04": "ffl",
            "\ufb06": "st",
            "\u017ft": "ft",
            "\u2014": "--",
            "\u2013": "-",
            "\u2018": "'",
            "\u2019": "'",
            "\u201c": '"',
            "\u201d": '"',
            "\u2026": "...",
        }

        for ligature, replacement in ligature_map.items():
            if ligature in text:
                text = text.replace(ligature, replacement)

        return text

    def _normalize_utf8(self, text: str) -> str:
        """Robust UTF-8 normalization and mojibake fixing"""
        if not text:
            return text

        # Common mojibake patterns (double-encoded UTF-8)
        mojibake_fixes = {
            # German umlauts
            "ÃƒÂ¤": "ä",
            "ÃƒÂ¶": "ö",
            "ÃƒÂ¼": "ü",
            "Ãƒâ€ž": "Ä",
            'Ãƒâ€"': "Ö",
            'ÃƒÅ"': "Ü",
            "ÃƒÅ¸": "ß",
            # French/Spanish
            "ÃƒÂ©": "é",
            "ÃƒÂ¨": "è",
            "ÃƒÂ¡": "á",
            "Ãƒ ": "à",
            "ÃƒÂ³": "ó",
            "ÃƒÂ²": "ò",
            "ÃƒÂº": "ú",
            "ÃƒÂ¹": "ù",
            "ÃƒÂ­": "í",
            "ÃƒÂ¬": "ì",
            "ÃƒÂ±": "ñ",
            "ÃƒÂ§": "ç",
            # Quotes and dashes
            'Ã¢â‚¬Å"': '"',
            "Ã¢â‚¬": '"',
            "Ã¢â‚¬â„¢": "'",
            "Ã¢â‚¬Ëœ": "'",
            'Ã¢â‚¬"': "–",
            'Ã¢â‚¬"': "—",
            "Ã¢â‚¬Â¦": "...",
            "Ã¢â‚¬Â¢": "•",
        }

        # Apply mojibake fixes
        for wrong, correct in mojibake_fixes.items():
            text = text.replace(wrong, correct)

        # Normalize to NFC form (canonical composition)
        try:
            text = unicodedata.normalize("NFC", text)
        except Exception:
            pass

        return text

    def _fix_text_formatting_issues(self, text: str) -> str:
        """Post-processing to fix hyphenation and line breaks"""
        # Remove PDF artifacts
        text = text.replace('Ã¢â‚¬"', "")

        # Fix hyphenated words
        def fix_hyphenation(match):
            word1 = match.group(1)
            word2 = match.group(2)
            common_words = {
                "und",
                "oder",
                "aber",
                "sowie",
                "als",
                "auf",
                "aus",
                "bei",
                "der",
                "die",
                "das",
                "ein",
                "eine",
                "mit",
                "von",
                "zu",
                "in",
                "an",
                "für",
                "über",
                "unter",
                "and",
                "or",
                "but",
                "the",
                "a",
                "an",
                "at",
                "on",
                "in",
                "for",
                "with",
            }
            if word2.lower() in common_words:
                return word1 + " " + word2
            else:
                return word1 + word2

        text = re.sub(r"(\w{2,})-\s*\n\s*([a-zäöüß]\w*)", fix_hyphenation, text)
        text = re.sub(r"(\w{3,})-\s+([a-zäöüß]{4,})", r"\1\2", text)

        # Cleanup
        text = REGEX_PATTERNS["multiple_spaces"].sub(" ", text)
        text = re.sub(r"\s+([.,;:!?])", r"\1", text)
        text = re.sub(r"([.,;:!?])(?=[A-ZÄÖÜäöüß])", r"\1 ", text)
        text = re.sub(r"(\w)(\d{1,2})([A-ZÄÖÜ])", r"\1\2 \3", text)
        text = REGEX_PATTERNS["excessive_newlines"].sub("\n\n\n", text)

        return text

    def extract_text(self) -> str:
        """
        Main extraction method with advanced layout recognition.
        Column detection and adaptive columns are always enabled.
        """
        print(f"Processing PDF: {self.pdf_path}")
        print(f"Total pages: {len(self.doc)}")

        # Extract all blocks with positions
        all_blocks = self._extract_all_blocks_with_positions()

        # Apply intelligent layout analysis and sorting (always enabled)
        all_blocks = self._apply_intelligent_sorting(all_blocks)

        # Build text from sorted blocks
        final_text = self._build_text_from_blocks(all_blocks)

        self.doc.close()
        return final_text

    def _build_text_from_blocks(self, blocks: List[TextBlock]) -> str:
        """Build complete text from sorted blocks in correct reading order"""
        if not blocks:
            return ""

        pages_blocks = {}
        for block in blocks:
            if block.page_num not in pages_blocks:
                pages_blocks[block.page_num] = []
            pages_blocks[block.page_num].append(block)

        text_parts = []

        for page_num in sorted(pages_blocks.keys()):
            page_blocks = pages_blocks[page_num]

            logger.debug(f"Building text for page {page_num + 1} with {len(page_blocks)} blocks")

            if page_num == 0:
                text_parts.append(f"--- Seite {page_num + 1} ---\n")
            else:
                text_parts.append(f"\n\n--- Seite {page_num + 1} ---\n")

            current_column = None
            last_block = None
            paragraph_buffer = []
            last_was_special = False

            for block in page_blocks:
                if block.semantic_type in [
                    BlockType.HEADER,
                    BlockType.FOOTER,
                    BlockType.PAGE_NUMBER,
                ]:
                    continue

                if block.column_id != current_column and block.column_id is not None:
                    if paragraph_buffer:
                        text_parts.append(" ".join(paragraph_buffer))
                        paragraph_buffer = []

                    current_column = block.column_id
                    last_block = None

                    if text_parts and not text_parts[-1].endswith("\n\n"):
                        text_parts.append("\n\n")

                needs_paragraph_break = False
                needs_space = True

                if last_block:
                    y_gap = abs(self._get_sort_y(block) - self._get_sort_y(last_block))
                    x_diff = abs(block.bbox[0] - last_block.bbox[0])

                    if block.semantic_type != last_block.semantic_type:
                        needs_paragraph_break = True
                    elif y_gap > last_block.height * 2.0:
                        needs_paragraph_break = True
                    elif x_diff > 50 and block.semantic_type == BlockType.MAIN_TEXT:
                        needs_paragraph_break = True
                    elif y_gap < last_block.height * 1.5:
                        needs_space = True
                        needs_paragraph_break = False

                if block.semantic_type == BlockType.TABLE:
                    if paragraph_buffer:
                        text_parts.append(" ".join(paragraph_buffer))
                        paragraph_buffer = []
                    if not last_was_special:
                        text_parts.append("\n")
                    text_parts.append(block.text)
                    text_parts.append("\n")
                    last_was_special = True

                elif block.semantic_type == BlockType.FOOTNOTE:
                    if paragraph_buffer:
                        text_parts.append(" ".join(paragraph_buffer))
                        paragraph_buffer = []
                    text_parts.append(f"\n{block.text}")
                    last_was_special = True

                elif block.semantic_type == BlockType.CAPTION:
                    if paragraph_buffer:
                        text_parts.append(" ".join(paragraph_buffer))
                        text_parts.append(" ")
                    text_parts.append(block.text)
                    last_was_special = False

                elif block.semantic_type in [
                    BlockType.INFOBOX,
                    BlockType.MARGIN_LEFT,
                    BlockType.MARGIN_RIGHT,
                ]:
                    if paragraph_buffer:
                        if (
                            last_block
                            and self._get_sort_y(block)
                            > self._get_sort_y(last_block) + last_block.height * 2
                        ):
                            text_parts.append(" ".join(paragraph_buffer))
                            text_parts.append("\n\n")
                            text_parts.append(block.text)
                            paragraph_buffer = []
                        else:
                            paragraph_buffer.append(block.text)
                    else:
                        paragraph_buffer = [block.text]
                    last_was_special = False

                else:
                    if last_was_special:
                        if paragraph_buffer:
                            text_parts.append(" ".join(paragraph_buffer))
                        text_parts.append("\n\n")
                        paragraph_buffer = [block.text]
                    elif needs_paragraph_break and paragraph_buffer:
                        text_parts.append(" ".join(paragraph_buffer))
                        text_parts.append("\n\n")
                        paragraph_buffer = [block.text]
                    else:
                        if paragraph_buffer and needs_space:
                            if paragraph_buffer[-1].endswith("-"):
                                paragraph_buffer[-1] = paragraph_buffer[-1][:-1] + block.text
                            else:
                                paragraph_buffer.append(block.text)
                        else:
                            paragraph_buffer.append(block.text)
                    last_was_special = False

                last_block = block

            if paragraph_buffer:
                text_parts.append(" ".join(paragraph_buffer))

        result = "\n".join(text_parts)

        # Clean up whitespace
        result = REGEX_PATTERNS["excessive_newlines"].sub("\n\n\n", result)
        result = REGEX_PATTERNS["triple_newlines"].sub("\n\n", result)
        result = REGEX_PATTERNS["multiple_spaces"].sub(" ", result)
        result = REGEX_PATTERNS["space_after_newline"].sub("\n", result)
        result = REGEX_PATTERNS["space_before_newline"].sub("\n", result)
        result = result.strip()

        # Fix text formatting and normalize UTF-8
        result = self._fix_text_formatting_issues(result)
        result = self._normalize_utf8(result)

        return result

    def _extract_all_blocks_with_positions(self) -> List[TextBlock]:
        """Extract ALL blocks with complete position information"""
        all_blocks = []

        for page_num, page in enumerate(self.doc):
            page_layout = self.page_layouts[page_num]

            if page_num % 10 == 0:
                logger.debug(f"Extracting blocks from page {page_num + 1}/{len(self.doc)}...")

            blocks_dict = page.get_text("dict")

            for block in blocks_dict.get("blocks", []):
                if block.get("type") == 0:
                    for line in block.get("lines", []):
                        line_spans = []
                        line_bboxes = []

                        for span in line.get("spans", []):
                            span_text = span.get("text", "")
                            span_text = self._fix_ligatures(span_text)
                            span_text = self._normalize_utf8(span_text)

                            if span_text.strip():
                                line_spans.append(span_text)
                                line_bboxes.append(span.get("bbox", (0, 0, 0, 0)))

                        if line_spans:
                            combined_text = ""
                            for i, span_text in enumerate(line_spans):
                                if i == 0:
                                    combined_text = span_text
                                else:
                                    prev_bbox = line_bboxes[i - 1]
                                    curr_bbox = line_bboxes[i]
                                    gap = curr_bbox[0] - prev_bbox[2]

                                    if gap < 2:
                                        combined_text += span_text
                                    else:
                                        combined_text += " " + span_text

                            text = combined_text.strip()
                            if text:
                                combined_bbox = (
                                    min(bbox[0] for bbox in line_bboxes),
                                    min(bbox[1] for bbox in line_bboxes),
                                    max(bbox[2] for bbox in line_bboxes),
                                    max(bbox[3] for bbox in line_bboxes),
                                )

                                font_size = line.get("spans", [{}])[0].get("size", 0)

                                text_block = TextBlock(
                                    text=text,
                                    bbox=combined_bbox,
                                    page_num=page_num,
                                    block_type="text",
                                    font_size=font_size,
                                )

                                text_block.semantic_type = self._classify_block_type(
                                    text_block, page_layout
                                )

                                all_blocks.append(text_block)

            try:
                tables = self._extract_tables_from_page(page, page_num)
                all_blocks.extend(tables)
            except Exception as e:
                logger.debug(f"Table extraction error on page {page_num}: {e}")

        return all_blocks

    def _classify_block_type(self, block: TextBlock, page_layout: Dict) -> BlockType:
        """Classify the semantic type of a text block"""
        page_width = page_layout["width"]
        page_height = page_layout["height"]

        x_ratio = block.x_center / page_width
        y_ratio = block.y_center / page_height

        text = block.text.strip()
        if block.width < 100 and block.height < 50:
            if y_ratio < 0.15 or y_ratio > 0.85:
                if REGEX_PATTERNS["page_number"].match(text) or REGEX_PATTERNS["page_text"].match(
                    text
                ):
                    return BlockType.PAGE_NUMBER

        if block.font_size:
            avg_font_size = 11
            is_title_candidate = False

            if block.font_size > avg_font_size * 1.3:
                is_title_candidate = True

            if 0.3 < x_ratio < 0.7 and y_ratio < 0.3:
                is_title_candidate = True

            if len(text) < 100 and (text[0].isupper() if text else False):
                if re.match(r"^\d+\.?\s+[A-Z]", text) or re.match(
                    r"^(Chapter|Section|Part|Teil|Kapitel)\s+\d+", text, re.IGNORECASE
                ):
                    is_title_candidate = True

            if is_title_candidate:
                return BlockType.MAIN_TEXT

        if y_ratio < 0.08 and block.font_size and block.font_size < 10:
            return BlockType.HEADER
        elif y_ratio > 0.92 and block.font_size and block.font_size < 10:
            return BlockType.FOOTER

        if block.width < page_width * 0.2:
            if x_ratio < 0.15:
                return BlockType.MARGIN_LEFT
            elif x_ratio > 0.85:
                return BlockType.MARGIN_RIGHT

        if block.font_size and block.font_size < 9 and y_ratio > 0.7:
            return BlockType.FOOTNOTE

        if block.font_size and block.font_size < 10 and "Figure" in block.text:
            return BlockType.CAPTION

        if block.width < page_width * 0.35 and block.width > page_width * 0.15:
            return BlockType.INFOBOX

        return BlockType.MAIN_TEXT

    def _apply_intelligent_sorting(self, blocks: List[TextBlock]) -> List[TextBlock]:
        """Apply intelligent sorting with advanced column detection"""
        if not blocks:
            return blocks

        pages_blocks = {}
        for block in blocks:
            if block.page_num not in pages_blocks:
                pages_blocks[block.page_num] = []
            pages_blocks[block.page_num].append(block)

        ordered_blocks = []

        for page_num in sorted(pages_blocks.keys()):
            page_blocks = pages_blocks[page_num]
            page_layout = self.page_layouts[page_num]

            logger.debug(f"Processing page {page_num + 1} with {len(page_blocks)} blocks")

            if page_layout["layout_type"] == LayoutType.LANDSCAPE_OR_DOUBLE:
                ordered = self._process_landscape_page(page_blocks, page_layout)
            else:
                ordered = self._process_single_page(page_blocks, page_layout)

            ordered_blocks.extend(ordered)

        return ordered_blocks

    def _process_landscape_page(
        self, blocks: List[TextBlock], page_layout: Dict
    ) -> List[TextBlock]:
        """Process landscape page with improved double page detection"""
        mid_x = page_layout["width"] / 2
        page_width = page_layout["width"]
        page_height = page_layout["height"]

        gap_empty_ratio = self._analyze_middle_gap_coverage(blocks, mid_x)

        left_blocks = [b for b in blocks if b.x_center < mid_x]
        right_blocks = [b for b in blocks if b.x_center >= mid_x]

        total_blocks = len(blocks)
        if total_blocks > 0:
            balance_ratio = len(left_blocks) / total_blocks
        else:
            balance_ratio = 0.5

        min_blocks_per_side = max(3, total_blocks * 0.1)

        page_number_info = self._detect_page_numbers(blocks, page_width, page_height)

        confidence = 0.0

        if gap_empty_ratio > 0.8:
            confidence += 40
        elif gap_empty_ratio > 0.5:
            confidence += 20

        if 0.25 < balance_ratio < 0.75:
            confidence += 30
        elif 0.15 < balance_ratio < 0.85:
            confidence += 15

        if len(left_blocks) >= min_blocks_per_side and len(right_blocks) >= min_blocks_per_side:
            confidence += 20

        if page_number_info["left"] or page_number_info["right"]:
            if not page_number_info["center"]:
                confidence += 10
        elif page_number_info["center"]:
            confidence -= 20

        is_double_page = confidence >= 50

        logger.debug(
            f"Landscape page {page_layout['page_num'] + 1}: confidence={confidence}/100, {'DOUBLE' if is_double_page else 'SINGLE'}"
        )

        if is_double_page:
            left_ordered = self._detect_and_order_columns_advanced(left_blocks, page_width / 2)
            right_ordered = self._detect_and_order_columns_advanced(right_blocks, page_width / 2)

            if left_ordered and right_ordered:
                avg_left_x = sum(b.x_center for b in left_ordered) / len(left_ordered)
                avg_right_x = sum(b.x_center for b in right_ordered) / len(right_ordered)

                if avg_left_x > avg_right_x:
                    logger.debug("Pages seem reversed, swapping...")
                    return right_ordered + left_ordered
                else:
                    return left_ordered + right_ordered
            else:
                return left_ordered + right_ordered
        else:
            return self._detect_and_order_columns_advanced(blocks, page_width)

    def _process_single_page(self, blocks: List[TextBlock], page_layout: Dict) -> List[TextBlock]:
        """Process single page with flexible column detection"""
        return self._detect_and_order_columns_advanced(blocks, page_layout["width"])

    def _detect_and_order_columns_advanced(
        self, blocks: List[TextBlock], page_width: float
    ) -> List[TextBlock]:
        """Advanced column detection without artificial limits"""
        if len(blocks) < 2:
            return blocks

        gaps = self._find_all_vertical_gaps(blocks, page_width)
        logger.debug(f"Found {len(gaps)} column boundaries")

        if not gaps:
            blocks.sort(key=lambda b: (self._get_sort_y(b), b.bbox[0]))
            return blocks

        columns = self._assign_blocks_to_columns(blocks, gaps)
        valid_columns = self._validate_columns(columns)

        columns_with_position = []
        for col_blocks in valid_columns:
            if col_blocks:
                avg_x = sum(b.x_center for b in col_blocks) / len(col_blocks)
                columns_with_position.append((avg_x, col_blocks))

        columns_with_position.sort(key=lambda x: x[0])

        ordered = []
        for col_idx, (avg_x, col_blocks) in enumerate(columns_with_position):
            for block in col_blocks:
                block.column_id = col_idx

            col_blocks.sort(key=lambda b: (self._get_sort_y(b), b.bbox[0]))
            ordered.extend(col_blocks)

        return ordered

    def _find_all_vertical_gaps(self, blocks: List[TextBlock], page_width: float) -> List[float]:
        """Find all significant vertical gaps"""
        if not blocks:
            return []

        coverage = {}
        for block in blocks:
            for x in range(int(block.bbox[0]), int(block.bbox[2])):
                coverage[x] = coverage.get(x, 0) + 1

        if not coverage:
            return []

        gaps = []
        gap_start = None
        for x in sorted(coverage.keys()):
            if coverage.get(x, 0) == 0:
                if gap_start is None:
                    gap_start = x
            else:
                if gap_start is not None and (x - gap_start) >= MIN_GAP_WIDTH:
                    gap_center = (gap_start + x) / 2
                    gaps.append(gap_center)
                gap_start = None

        if len(gaps) < 2:
            x_positions = [block.bbox[0] for block in blocks]
            cluster_gaps = self._find_gaps_by_clustering(x_positions)
            gaps.extend(cluster_gaps)
            gaps = sorted(list(set(gaps)))

        return sorted(gaps)

    def _find_gaps_by_clustering(self, x_positions: List[float]) -> List[float]:
        """Find gaps using clustering of x-coordinates"""
        if len(x_positions) < 4:
            return []

        sorted_x = sorted(x_positions)
        gaps = []

        if len(sorted_x) > 2:
            diffs = [sorted_x[i] - sorted_x[i - 1] for i in range(1, len(sorted_x))]
            mean_diff = statistics.mean(diffs)
            std_diff = statistics.stdev(diffs) if len(diffs) > 1 else mean_diff
            threshold = mean_diff + std_diff * 1.5
        else:
            threshold = 50

        for i in range(1, len(sorted_x)):
            gap_size = sorted_x[i] - sorted_x[i - 1]
            if gap_size > threshold:
                boundary = (sorted_x[i - 1] + sorted_x[i]) / 2
                gaps.append(boundary)

        return gaps

    def _assign_blocks_to_columns(
        self, blocks: List[TextBlock], gaps: List[float]
    ) -> List[List[TextBlock]]:
        """Assign blocks to columns based on gap positions"""
        if not gaps:
            return [blocks]

        columns = [[] for _ in range(len(gaps) + 1)]

        for block in blocks:
            x_center = block.x_center
            col_idx = 0

            for gap in gaps:
                if x_center > gap:
                    col_idx += 1
                else:
                    break

            if col_idx < len(columns):
                columns[col_idx].append(block)

        return columns

    def _validate_columns(self, columns: List[List[TextBlock]]) -> List[List[TextBlock]]:
        """Validate columns - remove empty or too narrow ones"""
        valid = []

        for col in columns:
            if not col:
                continue

            min_x = min(b.bbox[0] for b in col)
            max_x = max(b.bbox[2] for b in col)
            width = max_x - min_x

            if width > 30 or len(col) > 2:
                valid.append(col)

        return valid

    def _extract_tables_from_page(self, page, page_num: int) -> List[TextBlock]:
        """Extract tables from a page"""
        table_blocks = []

        try:
            tables = page.find_tables()
            for table in tables:
                table_text = self._format_table(table)
                if table_text:
                    table_block = TextBlock(
                        text=table_text,
                        bbox=table.bbox,
                        page_num=page_num,
                        block_type="table",
                        semantic_type=BlockType.TABLE,
                    )
                    table_blocks.append(table_block)
        except Exception as e:
            logger.debug(f"Table detection error on page {page_num}: {e}")

        return table_blocks

    def _format_table(self, table) -> str:
        """Format table data as markdown"""
        try:
            data = table.extract()
            if not data:
                return ""

            md_lines = []
            for i, row in enumerate(data):
                if row:
                    row_text = " | ".join(str(cell or "").strip() for cell in row)
                    if row_text.strip():
                        md_lines.append(f"| {row_text} |")

                    if i == 0 and len(md_lines) > 0:
                        separator = " | ".join("---" for _ in row)
                        md_lines.append(f"| {separator} |")

            return "\n".join(md_lines) if md_lines else ""

        except Exception as e:
            logger.debug(f"Table formatting error: {e}")
            return ""


# ============= DOCUMENT PROCESSOR =============
@dataclass
class ExtractionResult:
    """Container for extraction results with quality metrics"""

    text: str
    metadata: Dict = field(default_factory=dict)
    quality_score: float = 0.0
    extraction_method: str = "unknown"
    tables: List = field(default_factory=list)
    processing_time: float = 0.0
    warnings: List = field(default_factory=list)


class DocumentProcessor:
    """Streamlined document processor for GraphRAG pipeline"""

    def __init__(
        self,
        config_path: str = "config.yaml",
        config_section: str = "document_processing",
        **kwargs,
    ):
        """
        Initialize with YAML config and optional overrides

        Args:
            config_path: Path to YAML configuration file
            config_section: Section in config file containing document processing settings
            **kwargs: Override any config values (including track_performance)
        """
        self.config = self._load_config(config_path, config_section)

        # Allow override of specific config values including track_performance
        for key, value in kwargs.items():
            if key in [
                "output_dir",
                "preserve_filenames",
                "enable_ocr",
                "quality_threshold",
                "force_reprocess",
                "track_performance",
            ]:
                self.config[key] = value

        self.output_dir = Path(self.config["output_dir"])
        self.preserve_filenames = self.config["preserve_filenames"]
        self.enable_ocr = self.config["enable_ocr"]
        self.quality_threshold = self.config["quality_threshold"]
        self.force_reprocess = self.config.get("force_reprocess", False)
        self.supported_formats = set(self.config["supported_formats"])

        # Performance tracking configuration (default: disabled)
        self.track_performance = self.config.get("track_performance", False)
        self.save_performance_data = self.config["performance"].get("save_performance_data", False)

        # Create output directories
        self.output_dir.mkdir(exist_ok=True, parents=True)
        for subdir in ["texts", "metadata"]:
            (self.output_dir / subdir).mkdir(exist_ok=True)

        # Only create performance folder if saving performance data is enabled
        if self.save_performance_data:
            (self.output_dir / "performance").mkdir(exist_ok=True)

        # Performance tracking stats (only if enabled)
        if self.track_performance:
            self.processing_stats = {
                "total_documents": 0,
                "total_processing_time": 0.0,
                "extraction_methods": {},
                "quality_scores": [],
            }

    def _load_config(self, config_path: str, config_section: str) -> Dict:
        """Load configuration from YAML file"""
        default_config = {
            "output_dir": "../data/step1_processed",
            "preserve_filenames": True,
            "quality_threshold": 0.5,
            "enable_ocr": False,
            "force_reprocess": False,
            "track_performance": False,
            "extraction": {
                "max_text_sample": 5000,
                "min_text_for_ocr": 50,
                "scanned_page_threshold": 0.5,
                "text_density_threshold": 1000,
                "max_sample_pages": 10,
            },
            "quality_scoring": {"chars_per_page_norm": 1000, "artifact_penalty_divisor": 100},
            "supported_formats": [".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md"],
            "performance": {
                "track_performance": False,
                "save_performance_data": False,
                "log_warnings": True,
            },
        }

        config_file = Path(config_path)
        if config_file.exists():
            try:
                with open(config_file, "r") as f:
                    full_config = yaml.safe_load(f) or {}

                if config_section in full_config:
                    loaded = full_config[config_section]

                    for key, value in loaded.items():
                        if isinstance(value, dict) and key in default_config:
                            if isinstance(default_config[key], dict):
                                default_config[key].update(value)
                            else:
                                default_config[key] = value
                        else:
                            default_config[key] = value
                else:
                    print(
                        f"Warning: Section '{config_section}' not found in {config_path}, using defaults"
                    )

            except Exception as e:
                print(f"Using default config. Error loading {config_path}: {e}")
        else:
            print(f"Config file {config_path} not found, using defaults")

        return default_config

    def generate_doc_id(self, file_path: str) -> str:
        """Generate unique document ID with content hash"""
        file_path = Path(file_path)

        try:
            with open(file_path, "rb") as f:
                content = f.read()
                content_hash = hashlib.md5(content).hexdigest()[:8]
        except Exception:
            content_hash = hashlib.md5(str(file_path).encode()).hexdigest()[:8]

        if self.preserve_filenames:
            clean_name = re.sub(r"[^\w\-_]", "_", file_path.stem)[:50]
            return f"{clean_name}_{content_hash}"

        return content_hash

    def calculate_extraction_quality(self, text: str, page_count: int = 1) -> float:
        """Calculate extraction quality score"""
        if not text or page_count == 0:
            return 0.0

        config = self.config["quality_scoring"]
        sample = text[: self.config["extraction"]["max_text_sample"]]

        density = min(len(text) / (page_count * config["chars_per_page_norm"]), 1.0)
        readable = sum(c.isalnum() or c.isspace() for c in sample) / len(sample) if sample else 0
        no_artifacts = 1.0 - (
            len(re.findall(r"[\x00-\x1f]", sample)) / config["artifact_penalty_divisor"]
        )

        return round((density + readable + no_artifacts) / 3, 3)

    def extract_pdf(self, file_path: str) -> ExtractionResult:
        """Extract text from PDF using advanced hybrid layout recognition"""
        start_time = time.time()
        warnings = []
        extractor = None

        try:
            try:
                import pymupdf4llm
            except ImportError:
                warnings.append("pymupdf4llm not available - required for PDF extraction")
                return ExtractionResult(
                    text="",
                    quality_score=0.0,
                    extraction_method="failed",
                    processing_time=time.time() - start_time,
                    warnings=warnings,
                )

            print("    Extracting with hybrid approach (pymupdf4llm + advanced layout analysis)...")

            extractor = PDFComplexLayoutExtractor(pdf_path=str(file_path), debug=False)

            extracted_text = extractor.extract_text()

            page_count = len(extractor.page_layouts)
            table_count = extracted_text.count("|---")

            metadata = {
                "page_count": page_count,
                "extraction_enhanced": True,
                "layout_preserved": True,
                "extraction_method": "hybrid_advanced_layout",
                "pages": [
                    {
                        "page_number": i + 1,
                        "layout_type": layout["layout_type"].value,
                        "width": layout["width"],
                        "height": layout["height"],
                    }
                    for i, layout in enumerate(extractor.page_layouts)
                ],
            }

            if table_count > 0:
                metadata["has_tables"] = True
                metadata["table_count"] = table_count
                warnings.append(f"{table_count} tables detected and preserved in markdown format")

            quality = self.calculate_extraction_quality(extracted_text, page_count)
            print(f"    Quality score: {quality:.2f}")

            if quality < 0.1:
                warnings.append(f"Very low quality extraction ({quality:.2f})")

            print(f"    Extracted {len(extracted_text)} characters from {page_count} pages")
            print("    Successfully extracted with advanced hybrid approach")

            result = ExtractionResult(
                text=extracted_text,
                metadata=metadata,
                quality_score=quality,
                extraction_method="hybrid_advanced_layout",
                processing_time=time.time() - start_time,
                warnings=warnings,
            )

            if extractor and hasattr(extractor, "doc"):
                try:
                    extractor.doc.close()
                except:
                    pass

            return result

        except Exception as e:
            print(f"    PDF extraction failed: {str(e)[:100]}")
            warnings.append(f"PDF extraction failed: {str(e)[:100]}")

            if extractor and hasattr(extractor, "doc"):
                try:
                    extractor.doc.close()
                except:
                    pass

            return ExtractionResult(
                text="",
                quality_score=0.0,
                extraction_method="failed",
                processing_time=time.time() - start_time,
                warnings=warnings,
            )

    def extract_docx(self, file_path: str) -> ExtractionResult:
        """Extract text from DOCX"""
        start_time = time.time()
        warnings = []

        try:
            doc = DocxDocument(file_path)
            text = ""
            paragraph_count = 0

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n\n"
                    paragraph_count += 1

            table_count = len(doc.tables)
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    if any(row_text):
                        text += " | ".join(row_text) + "\n"
                text += "\n"

            props = doc.core_properties
            metadata = {
                "paragraph_count": paragraph_count,
                "table_count": table_count,
                "page_count": 0,
                "docx_metadata": {
                    "author": props.author or "",
                    "created": props.created.isoformat() if props.created else "",
                    "modified": props.modified.isoformat() if props.modified else "",
                },
            }

            if table_count > 0:
                metadata["has_tables"] = True
                warnings.append(f"Document contains {table_count} tables")

            return ExtractionResult(
                text=text,
                metadata=metadata,
                quality_score=self.calculate_extraction_quality(text),
                extraction_method="python-docx",
                processing_time=time.time() - start_time,
                warnings=warnings,
            )
        except Exception as e:
            warnings.append(f"DOCX extraction failed: {str(e)[:100]}")
            return ExtractionResult(
                text="",
                quality_score=0.0,
                extraction_method="failed",
                processing_time=time.time() - start_time,
                warnings=warnings,
            )

    def extract_pptx(self, file_path: str) -> ExtractionResult:
        """Extract text from PPTX"""
        start_time = time.time()
        warnings = []

        try:
            prs = Presentation(file_path)
            text = ""
            slides_data = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = ""

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"

                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            if any(row_text):
                                slide_text += " | ".join(row_text) + "\n"

                if slide_text.strip():
                    slides_data.append({"slide_number": slide_num, "char_count": len(slide_text)})
                    text += f"\n--- Slide {slide_num} ---\n{slide_text}"

            metadata = {
                "slide_count": len(prs.slides),
                "page_count": len(slides_data),
                "slides": slides_data,
            }

            return ExtractionResult(
                text=text,
                metadata=metadata,
                quality_score=self.calculate_extraction_quality(text),
                extraction_method="python-pptx",
                processing_time=time.time() - start_time,
                warnings=warnings,
            )
        except Exception as e:
            warnings.append(f"PPTX extraction failed: {str(e)[:100]}")
            return ExtractionResult(
                text="",
                quality_score=0.0,
                extraction_method="failed",
                processing_time=time.time() - start_time,
                warnings=warnings,
            )

    def extract_xlsx(self, file_path: str) -> ExtractionResult:
        """Extract text from XLSX"""
        start_time = time.time()
        warnings = []

        try:
            workbook = load_workbook(file_path, data_only=True)
            text = ""
            sheets_data = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = f"\n--- Sheet: {sheet_name} ---\n"
                non_empty_rows = 0

                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(v) if v is not None else "" for v in row]
                    row_text = " | ".join(row_values).strip()
                    if row_text:
                        sheet_text += row_text + "\n"
                        non_empty_rows += 1

                if non_empty_rows > 0:
                    sheets_data.append(
                        {
                            "sheet_name": sheet_name,
                            "rows": non_empty_rows,
                            "char_count": len(sheet_text),
                        }
                    )
                    text += sheet_text

            metadata = {
                "sheet_count": len(workbook.sheetnames),
                "sheets": sheets_data,
                "page_count": 0,
            }
            workbook.close()

            return ExtractionResult(
                text=text,
                metadata=metadata,
                quality_score=self.calculate_extraction_quality(text),
                extraction_method="openpyxl",
                processing_time=time.time() - start_time,
                warnings=warnings,
            )
        except Exception as e:
            warnings.append(f"XLSX extraction failed: {str(e)[:100]}")
            return ExtractionResult(
                text="",
                quality_score=0.0,
                extraction_method="failed",
                processing_time=time.time() - start_time,
                warnings=warnings,
            )

    def extract_plaintext(self, file_path: str) -> ExtractionResult:
        """Extract text from plain text or markdown files"""
        start_time = time.time()
        warnings = []

        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

            is_markdown = Path(file_path).suffix.lower() == ".md"
            metadata = {
                "line_count": len(text.split("\n")),
                "is_markdown": is_markdown,
                "page_count": 0,
            }

            if is_markdown:
                headers = len(re.findall(r"^#{1,6}\s+", text, re.MULTILINE))
                code_blocks = len(re.findall(r"```[\s\S]*?```", text))
                metadata["markdown_stats"] = {
                    "header_count": headers,
                    "code_block_count": code_blocks,
                }

            return ExtractionResult(
                text=text,
                metadata=metadata,
                quality_score=self.calculate_extraction_quality(text),
                extraction_method="plain_text",
                processing_time=time.time() - start_time,
                warnings=warnings,
            )
        except Exception as e:
            warnings.append(f"Text extraction failed: {str(e)[:100]}")
            return ExtractionResult(
                text="",
                quality_score=0.0,
                extraction_method="failed",
                processing_time=time.time() - start_time,
                warnings=warnings,
            )

    def process_single_file(
        self, file_path: str, verbose: bool = True, force_reprocess: bool = None
    ) -> Optional[Dict]:
        """Process a single document file"""
        overall_start = time.time()
        file_path = Path(file_path)

        if force_reprocess is None:
            force_reprocess = self.force_reprocess

        if file_path.suffix.lower() not in self.supported_formats:
            if verbose:
                print(f"  Unsupported format: {file_path.suffix}")
            return None

        doc_id = self.generate_doc_id(str(file_path))
        text_file = self.output_dir / "texts" / f"{doc_id}.txt"

        extractors = {
            ".pdf": self.extract_pdf,
            ".docx": self.extract_docx,
            ".pptx": self.extract_pptx,
            ".xlsx": self.extract_xlsx,
            ".txt": self.extract_plaintext,
            ".md": self.extract_plaintext,
        }

        extractor = extractors.get(file_path.suffix.lower())
        if not extractor:
            return None

        try:
            result = extractor(str(file_path))

            if not result.text.strip():
                if verbose:
                    print(f"  No text content in {file_path.name}")
                return None

            if result.quality_score < self.quality_threshold:
                if verbose:
                    print(f"  Low quality ({result.quality_score:.2f}): {file_path.name}")

            total_time = time.time() - overall_start

            doc_data = {
                "doc_id": doc_id,
                "content_hash": doc_id.split("_")[-1],
                "original_filename": file_path.name,
                "filename": file_path.name,
                "file_path": str(file_path.absolute()),
                "file_type": file_path.suffix[1:].lower(),
                "file_size": os.path.getsize(file_path),
                "processed_at": datetime.now().isoformat(),
                "extraction_method": result.extraction_method,
                "extraction_quality": result.quality_score,
                "total_chars": len(result.text),
                "processing_time_seconds": round(total_time, 3),
                "extraction_time_seconds": round(result.processing_time, 3),
                **result.metadata,
            }

            if result.warnings:
                doc_data["extraction_warnings"] = result.warnings

            if result.tables:
                doc_data["table_count"] = len(result.tables)
                doc_data["has_tables"] = True

            with open(text_file, "w", encoding="utf-8") as f:
                f.write(result.text)

            meta_file = self.output_dir / "metadata" / f"{doc_id}.json"
            with open(meta_file, "w", encoding="utf-8") as f:
                json.dump(doc_data, f, indent=2, ensure_ascii=False)

            # Only save performance file if enabled
            if self.save_performance_data:
                perf_file = self.output_dir / "performance" / f"{doc_id}_performance.json"
                perf_data = {
                    "doc_id": doc_id,
                    "original_filename": file_path.name,
                    "processed_at": doc_data["processed_at"],
                    "extraction_quality": result.quality_score,
                    "extraction_method": result.extraction_method,
                    "extraction_time_seconds": round(result.processing_time, 3),
                    "total_processing_time_seconds": round(total_time, 3),
                    "chars_per_second": round(
                        (
                            len(result.text) / result.processing_time
                            if result.processing_time > 0
                            else 0
                        ),
                        0,
                    ),
                }
                with open(perf_file, "w", encoding="utf-8") as f:
                    json.dump(perf_data, f, indent=2, ensure_ascii=False)

            # Only track stats if enabled
            if self.track_performance:
                self.processing_stats["total_documents"] += 1
                self.processing_stats["total_processing_time"] += total_time
                self.processing_stats["quality_scores"].append(result.quality_score)

                if result.extraction_method not in self.processing_stats["extraction_methods"]:
                    self.processing_stats["extraction_methods"][result.extraction_method] = 0
                self.processing_stats["extraction_methods"][result.extraction_method] += 1

            if verbose:
                print(f"  {file_path.name} -> {doc_id}")
                print(
                    f"    Quality: {result.quality_score:.2f} | Method: {result.extraction_method} | "
                    f"Time: {total_time:.2f}s"
                )

            return doc_data

        except Exception as e:
            if verbose:
                print(f"  Failed: {file_path.name}: {e}")
            return None

    def process_directory(
        self,
        directory: str,
        patterns: List[str] = None,
        max_files: Optional[int] = None,
        force_reprocess: bool = None,
    ) -> List[Dict]:
        """Process all documents in a directory"""
        batch_start = time.time()
        directory = Path(directory)

        if force_reprocess is None:
            force_reprocess = self.force_reprocess

        if patterns is None:
            patterns = [f"*{ext}" for ext in self.supported_formats]

        files = []
        for pattern in patterns:
            files.extend(directory.glob(pattern))

        files = sorted(set(files))
        if max_files:
            files = files[:max_files]

        if not files:
            print("No supported files found")
            return []

        print(f"Found {len(files)} files to process")
        print(f"Output directory: {self.output_dir}\n")

        results = []
        failed_files = []
        skipped_count = 0

        for i, file_path in enumerate(files, 1):
            print(f"\n[{i}/{len(files)}] Processing: {file_path.name}")

            result = self.process_single_file(
                file_path, verbose=True, force_reprocess=force_reprocess
            )

            if result:
                results.append(result)
            elif (
                self.output_dir / "texts" / f"{self.generate_doc_id(str(file_path))}.txt"
            ).exists():
                skipped_count += 1
                doc_id = self.generate_doc_id(str(file_path))
                meta_file = self.output_dir / "metadata" / f"{doc_id}.json"
                if meta_file.exists():
                    with open(meta_file, "r", encoding="utf-8") as f:
                        results.append(json.load(f))
            else:
                failed_files.append(str(file_path))

        batch_time = time.time() - batch_start

        print(f"\n{'='*60}")
        print(" BATCH PROCESSING SUMMARY")
        print(f"{'='*60}")
        print(f"Success: {len(results)}/{len(files)} files")
        if skipped_count > 0:
            print(f"Skipped: {skipped_count} files (already processed)")
        if failed_files:
            print(f"Failed: {len(failed_files)} files")
            for file in failed_files[:3]:
                print(f"  - {Path(file).name}")
            if len(failed_files) > 3:
                print(f"  ... and {len(failed_files) - 3} more")

        # Only show stats if tracking enabled
        if results and self.track_performance and self.processing_stats.get("quality_scores"):
            avg_quality = sum(self.processing_stats["quality_scores"]) / len(
                self.processing_stats["quality_scores"]
            )
            print(f"Average quality: {avg_quality:.2f}")

            if self.processing_stats["extraction_methods"]:
                print("Extraction methods:")
                for method, count in self.processing_stats["extraction_methods"].items():
                    print(f"  - {method}: {count} files")

        print(f"Total time: {batch_time:.1f}s")
        print(f"Speed: {len(files)/batch_time:.1f} files/s")
        print(f"Output: {self.output_dir.absolute()}")
        print(f"{'='*60}")

        return results

    def list_processed_documents(self) -> List[Dict]:
        """List all already processed documents with their metadata"""
        metadata_folder = self.output_dir / "metadata"
        documents = []

        for meta_file in metadata_folder.glob("*.json"):
            try:
                with open(meta_file, "r", encoding="utf-8") as f:
                    documents.append(json.load(f))
            except Exception as e:
                print(f"Error loading metadata for {meta_file.name}: {e}")

        return documents

    def verify_processing_integrity(self) -> Dict:
        """Verify integrity of processed documents"""
        texts_folder = self.output_dir / "texts"
        metadata_folder = self.output_dir / "metadata"

        text_files = set(f.stem for f in texts_folder.glob("*.txt"))
        meta_files = set(f.stem for f in metadata_folder.glob("*.json"))

        missing_text = meta_files - text_files
        missing_meta = text_files - meta_files

        result = {
            "valid": len(missing_text) == 0 and len(missing_meta) == 0,
            "total_documents": len(text_files),
            "missing_text_files": list(missing_text),
            "missing_metadata_files": list(missing_meta),
        }

        if not result["valid"]:
            print("Integrity check failed:")
            if missing_text:
                print(f"  Missing text files: {len(missing_text)}")
            if missing_meta:
                print(f"  Missing metadata files: {len(missing_meta)}")
        else:
            print(f"Integrity check passed: {result['total_documents']} documents")

        return result


if __name__ == "__main__":
    processor = DocumentProcessor(config_path="config.yaml")
